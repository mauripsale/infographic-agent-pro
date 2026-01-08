from google.adk.agents import LlmAgent, SequentialAgent
from google.adk.tools import FunctionTool, GoogleSearchTool
import io
import os
import sys
import time
import json
from pathlib import Path
from bs4 import BeautifulSoup
import requests
from pptx import Presentation
from pptx.util import Inches

# Robust import for google.genai
try:
    from google import genai
    from google.genai import types
    HAS_GENAI = True
except ImportError:
    print("Warning: google.genai module not found. Image generation will be disabled.")
    HAS_GENAI = False

# Robust import for context
try:
    # Try importing context assuming it's in the python path (root)
    import context
    model_context = context.model_context
except ImportError:
    print("Warning: Could not import context. Using fallback ContextVar.")
    from contextvars import ContextVar
    model_context = ContextVar("model_context", default="gemini-2.5-flash")

STATIC_DIR = Path("static")

class DynamicLlmAgent(LlmAgent):
    """An agent that selects its model dynamically from context."""
    @property
    def model(self):
        return model_context.get()
    
    @model.setter
    def model(self, value):
        # Ignore setting model to fixed value, rely on context
        pass

def get_webpage_content(url: str) -> str:
    """Fetches and extracts the main text content from a URL."""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'lxml')
        # This is a simple heuristic, might need refinement
        main_content = soup.find('main') or soup.find('article') or soup.find('body')
        if main_content:
            return ' '.join(p.get_text() for p in main_content.find_all('p'))
        return "Could not extract main content."
    except requests.RequestException as e:
        return f"Error fetching URL: {e}"

def generate_image(prompt: str) -> io.BytesIO:
    """Generates an image using Gemini multimodal models exclusively."""
    if not HAS_GENAI:
        return None
    try:
        client = genai.Client(api_key=os.environ.get("GOOGLE_API_KEY"))
        
        # Determine which image model to use based on user selection
        current_model = model_context.get()
        
        # Mapping logic: if user is using a 'pro' model, use the pro image model.
        # Otherwise, default to the 2.5 flash image model.
        if "pro" in current_model.lower():
            model_to_use = "gemini-3-pro-image-preview"
        else:
            model_to_use = "gemini-2.5-flash-image"
            
        print(f"Generating image with model '{model_to_use}' for prompt: {prompt}")
        
        response = client.models.generate_content(
            model=model_to_use,
            contents=prompt
        )
        
        # Extract image bytes from the response candidates
        for part in response.candidates[0].content.parts:
            if part.inline_data:
                return io.BytesIO(part.inline_data.data)
            try:
                if hasattr(part, 'image_bytes') and part.image_bytes:
                    return io.BytesIO(part.image_bytes)
            except:
                pass
        
        return None
    except Exception as e:
        print(f"Image generation failed with {model_to_use}: {e}")
        return None

def create_presentation_file(json_content: str) -> str:
    """Creates a .pptx presentation file from a JSON structure."""
    try:
        slides_data = json.loads(json_content)
        prs = Presentation()
        
        for slide_info in slides_data:
            # Use a blank layout or Title/Content. 
            # Layout 1 is Title + Content. Layout 5 is Title Only. Layout 6 is Blank.
            # Let's use Layout 1 and resize content if image exists, or add image on side.
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = slide_info.get("title", "No Title")
            
            # Content (Bullet points)
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = ""
            for point in slide_info.get("bullet_points", []):
                p = content_placeholder.text_frame.add_paragraph()
                p.text = point
                p.level = 0 # Top level bullet
            
            # Image Generation
            image_prompt = slide_info.get("image_prompt")
            if image_prompt:
                print(f"Generating image for slide: {title.text}")
                image_stream = generate_image(image_prompt)
                if image_stream:
                    # Resize text box to left half, put image on right half.
                    content_placeholder.width = Inches(4.5)
                    # Add picture at Left=5.5 inches, Top=2 inches
                    slide.shapes.add_picture(image_stream, Inches(5), Inches(2), width=Inches(4.5))

        filename = STATIC_DIR / f"presentation_{int(time.time())}.pptx"
        prs.save(filename)
        return f"/static/{filename.name}"
    except Exception as e:
        # Log full error for debugging
        import traceback
        traceback.print_exc()
        return f"Error creating presentation file: {str(e)}"

# --- Agent 1: Script Generator ---
script_generator = DynamicLlmAgent(
    name="ScriptGenerator",
    model="gemini-2.5-flash", # Default, will be overridden by property
    description="Analyzes content from text or URLs and generates a presentation script.",
    instruction="""You are an expert content creator. Your task is to analyze the user's input (which can be plain text or a URL). If the input is a URL, use the `get_webpage_content` tool to fetch the text.
Synthesize the content into a structure for a slide presentation.
Generate a valid JSON output containing a list of slides. 
Each slide MUST have:
1. "title": A clear title.
2. "bullet_points": A list of strings (key points).
3. "image_prompt": A descriptive prompt for an AI image generator to create a visual relevant to the slide's content. Be creative and visual (e.g., "A futuristic city skyline with flying cars, neon lights, digital art style").

Output ONLY the JSON block. Do not add any conversational text before or after the JSON.""",
    tools=[FunctionTool(get_webpage_content)],
    output_key="slide_script"
)

# --- Agent 2: Slide Builder ---
slide_builder = DynamicLlmAgent(
    name="SlideBuilder",
    model="gemini-2.5-flash", # Default
    description="Generates a .pptx file from a presentation script.",
    instruction="""You are a presentation designer.
Take the JSON content provided in `slide_script` from the previous step.
Call the `create_presentation_file` tool with this JSON content.
Return the path to the generated .pptx file. Return ONLY the path, no other text.
""",
    tools=[FunctionTool(create_presentation_file)],
    output_key="pptx_file_path"
)

# --- Coordinator Agent: Sequential Pipeline ---
presentation_pipeline = SequentialAgent(
    name="InfographicAgent",
    description="A pipeline that transforms content into a presentation file.",
    sub_agents=[script_generator, slide_builder]
)
