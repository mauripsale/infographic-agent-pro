from google.adk.agents import LlmAgent, SequentialAgent
from google.adk.tools import FunctionTool, GoogleSearchTool
import json
import requests
from bs4 import BeautifulSoup
from pptx import Presentation

def get_webpage_content(url: str) -> str:
    """Fetches and extracts the main text content from a URL."""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'lxml')
        # This is a simple heuristic, might need refinement
        main_content = soup.find('main') or soup.find('article') or soup.find('body')
        if main_content:
            return ' '.join(p.get_text() for p in main_content.find_all('p'))
        return "Could not extract main content."
    except requests.RequestException as e:
        return f"Error fetching URL: {e}"

def create_presentation_file(json_content: str) -> str:
    """Creates a .pptx presentation file from a JSON structure."""
    try:
        slides_data = json.loads(json_content)
        prs = Presentation()
        for slide_info in slides_data:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_info.get("title", "No Title")
            
            content.text = ""
            for point in slide_info.get("bullet_points", []):
                p = content.text_frame.add_paragraph()
                p.text = point
                p.level = 1
        
        filename = "presentation_output.pptx"
        prs.save(filename)
        return f"Success! Presentation saved as {filename}"
    except Exception as e:
        return f"Error creating presentation file: {str(e)}"

# --- Agent 1: Script Generator ---
script_generator = LlmAgent(
    name="ScriptGenerator",
    model="gemini-1.5-flash-latest",
    description="Analyzes content from text or URLs and generates a presentation script.",
    instruction="""You are an expert content creator. Your task is to analyze the user's input (which can be plain text or a URL).
1. If the input is a URL, use the `get_webpage_content` tool to fetch the text.
2. Synthesize the content into a structure for a slide presentation.
3. Generate a valid JSON output containing a list of slides. Each slide must have "title" and "bullet_points" (a list of strings).
4. Output ONLY the JSON block. Do not add any conversational text before or after the JSON.""",
    tools=[FunctionTool(get_webpage_content)],
    output_key="slide_script"
)

# --- Agent 2: Slide Builder ---
slide_builder = LlmAgent(
    name="SlideBuilder",
    model="gemini-1.5-flash-latest",
    description="Generates a .pptx file from a presentation script.",
    instruction="""You are a presentation designer.
Take the JSON content provided in `slide_script` from the previous step.
Call the `create_presentation_file` tool with this JSON content.
Confirm to the user that the file is ready.""",
    tools=[FunctionTool(create_presentation_file)],
    output_key="final_result"
)

# --- Coordinator Agent: Sequential Pipeline ---
presentation_pipeline = SequentialAgent(
    name="InfographicAgent",
    description="A pipeline that transforms content into a presentation file.",
    sub_agents=[script_generator, slide_builder]
)
